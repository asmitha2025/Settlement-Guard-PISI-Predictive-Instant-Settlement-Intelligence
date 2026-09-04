import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Theme colors
    DARK_BG = RGBColor(11, 19, 43)        # Deep Navy
    CARD_BG = RGBColor(19, 30, 65)        # Lighter Card Navy
    CARD_BORDER = RGBColor(40, 60, 110)
    ACCENT_BLUE = RGBColor(56, 189, 248)  # Razorpay Sky Blue
    ACCENT_GREEN = RGBColor(52, 211, 153) # Mint/Success
    ACCENT_CORAL = RGBColor(248, 113, 113)# Alert Coral
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate Muted
    TEXT_YELLOW = RGBColor(251, 191, 36)
    ACCENT_YELLOW = RGBColor(251, 191, 36)

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title, category="TRACK 3 • AI REVENUE RECOVERY"):
        # Category tag
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        # Title
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)

    def add_card(slide, left, top, width, height, title=None, title_color=TEXT_WHITE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = title_color
        return card

    # ==========================================
    # SLIDE 1: Title
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s1)

    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "TRACK 3 • AI REVENUE RECOVERY"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p = tf1.add_paragraph()
    p.text = "PISI"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

    p = tf1.add_paragraph()
    p.text = "Predictive Instant Settlement Intelligence"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.space_before = Pt(4)

    p = tf1.add_paragraph()
    p.text = "Autonomous Post-Capture Liquidity Protection & Pre-Approved Instant Settlement"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(16)

    p = tf1.add_paragraph()
    p.text = "Razorpay AI Buildathon 2026"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_before = Pt(36)

    # ==========================================
    # SLIDE 2: The Problem
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s2)
    add_header(s2, "The Problem: Two Fundamentally Different Failures")

    # Left card
    add_card(s2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "1. AUTHORIZATION FAILURE", ACCENT_CORAL)
    tb = s2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets = [
        "The payment attempt doesn't go through.",
        "Customer's bank can't approve the debit.",
        "Razorpay error docs cite this as 'beyond our control' and point to multi-terminal routing.",
        "Handled by Smart Routing (Bygari et al., 2021).",
        "STATUS: ALREADY SOLVED BY RAZORPAY."
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED if i < 4 else ACCENT_CORAL
        p.font.bold = (i == 4)
        p.space_before = Pt(12)

    # Right card
    add_card(s2, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "2. SETTLEMENT RISK", ACCENT_GREEN)
    tb = s2.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets2 = [
        "The payment already succeeded and was captured.",
        "Funds are trapped in degraded partner bank CBS (T+2 to T+7 delay).",
        "Not a customer error — a severe merchant cash-flow crisis.",
        "Today, merchants must notice delays manually and pay 0.30%–0.50% fee.",
        "NOTHING PREDICTS IT AHEAD OF TIME. THIS IS WHAT PISI BUILDS."
    ]
    for i, b in enumerate(bullets2):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE if i < 4 else ACCENT_GREEN
        p.font.bold = (i == 4)
        p.space_before = Pt(12)

    # ==========================================
    # SLIDE 3: Razorpay Already Built Both Halves
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s3)
    add_header(s3, "Razorpay Already Built Both Halves")

    add_card(s3, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "SMART ROUTING", ACCENT_BLUE)
    tb = s3.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    sr_pts = [
        "94.69% Random Forest precision (best of 5 models).",
        "4–6% measured lift in production success rate.",
        "Trained on ~35 Million transactions.",
        "Operates exclusively on the acquiring/gateway side.",
        "Cannot touch settlement timing after capture."
    ]
    for i, pt in enumerate(sr_pts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(14)

    add_card(s3, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "INSTANT SETTLEMENT", ACCENT_YELLOW)
    tb = s3.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    is_pts = [
        "T+0 payout advance using Razorpay's corporate capital.",
        "Published fee at 0.30% per settlement.",
        "Applies only to payments that are already captured.",
        "100% merchant-triggered and reactive today.",
        "Only ~15% of merchants opt in; 85% remain unprotected."
    ]
    for i, pt in enumerate(is_pts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(14)

    # ==========================================
    # SLIDE 4: The Gap
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s4)
    add_header(s4, "The Gap: Nothing Connects Prediction to Capital")

    add_card(s4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.2))
    tb = s4.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SMART ROUTING                   INSTANT SETTLEMENT                  PISI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p = tf.add_paragraph()
    p.text = "Predicts which gateway          Advances payout,                    Predictive, automatic,"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

    p = tf.add_paragraph()
    p.text = "succeeds before auth            reactively upon request             pre-approved settlement"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED

    p = tf.add_paragraph()
    p.text = "(Acquiring side)                (Post-capture liquidity)            (Full Autonomous Loop)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    # Quote Callout
    callout = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(26, 46, 88)
    callout.line.color.rgb = ACCENT_BLUE
    tb = s4.shapes.add_textbox(Inches(1.1), Inches(5.2), Inches(11.1), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\"Smart Routing handles gateway-side degradation. Nothing predicts settlement-path risk ahead of time and pre-approves the advance automatically — today a merchant has to notice the delay themselves.\""
    p.font.size = Pt(17)
    p.font.italic = True
    p.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 5: Two Independent Legs
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s5)
    add_header(s5, "How It Works: Two Independent Legs")

    add_card(s5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "LEG A · SETTLEMENT PROTECTION", ACCENT_GREEN)
    tb = s5.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "MOVES CAPITAL (Autonomous Action)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    leg_a_pts = [
        "Monitors partner settlement banks via 5D Vitality.",
        "Predicts CBS delay 15–30 minutes before collapse.",
        "Pre-approves captured payments for T+0 credit.",
        "Discounted predictive fee: 0.10% (vs 0.30% reactive).",
        "Safety capped at 30% total portfolio allocation."
    ]
    for pt in leg_a_pts:
        p = tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(13.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(12)

    add_card(s5, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "LEG B · AUTHORIZATION WARNING", ACCENT_YELLOW)
    tb = s5.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "INFORMATIONAL ONLY (Zero Capital Deployed)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_YELLOW
    leg_b_pts = [
        "Detects customer issuing bank degradation patterns.",
        "Early warning webhooks push alerts to merchants.",
        "Allows merchant to steer checkout methods ahead of spikes.",
        "Never claims to prevent bank failure.",
        "Zero capital exposure; infinite scalability."
    ]
    for pt in leg_b_pts:
        p = tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(13.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(12)

    # ==========================================
    # SLIDE 6: Six-Layer Architecture
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s6)
    add_header(s6, "Architecture: Six Layers, Perceive to Learn")

    layers = [
        ("L1 Data Ingestion", "Captured-payment + failed-payment webhook streams (<5ms HMAC)", ACCENT_BLUE),
        ("L2 Feature Engineering", "5D Bank Vitality: Error acceleration, settlement velocity, maintenance windows", ACCENT_BLUE),
        ("L3 Prediction Engine", "Calibrated XGBoost downtime classifier (47 features) + duration estimate", ACCENT_GREEN),
        ("L4 Decision Engine", "3-Tier Escalation Matrix, 30% capital cap, 10 concurrent bridge limits", ACCENT_GREEN),
        ("L5 Execution Layer", "Instant Settlement execution + Immutable SHA-256 BridgeKeyID audit trail", ACCENT_YELLOW),
        ("L6 Monitoring & Drift", "Precision/recall tracking, Kolmogorov-Smirnov drift detection, auto-retrain", ACCENT_CORAL)
    ]
    top_pos = 1.4
    for title, desc, col in layers:
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = col
        card.line.width = Pt(1.2)

        tb = s6.shapes.add_textbox(Inches(1.1), Inches(top_pos + 0.1), Inches(3.2), Inches(0.65))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col

        tb2 = s6.shapes.add_textbox(Inches(4.4), Inches(top_pos + 0.1), Inches(7.8), Inches(0.65))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        top_pos += 0.95

    # ==========================================
    # SLIDE 7: Live System Console
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s7)
    add_header(s7, "Live System: Working Console & SHA-256 Audit Trail")

    add_card(s7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "BANK VITALITY SCORE", ACCENT_CORAL)
    tb = s7.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SBI Trajectory: 91 → 67 → 34 HP"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CORAL

    p = tf.add_paragraph()
    p.text = "Status: CRITICAL (Below 50 HP Activation Threshold)"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "Decisions Fired:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_before = Pt(20)

    p = tf.add_paragraph()
    p.text = "• Leg A: ACTIVATE (Autonomous Payout Advance)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = "• Leg B: WARN (Merchant Notification Pushed)"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_YELLOW
    p.space_before = Pt(6)

    add_card(s7, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "BRIDGE KEY ID (§7.3 SCHEMA)", ACCENT_BLUE)
    tb = s7.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sample Audit Record:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p = tf.add_paragraph()
    p.text = "Bridge ID: BRIDGE-SBI-20260823T133603-mo_000\nAmount: ₹3,641.63\nFee (0.10%): ₹3.64\nMerchant Credited: ₹3,637.99 (T+0 instant)"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = "SHA-256 Audit Hash (Web Crypto & Python hashlib):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_before = Pt(16)

    p = tf.add_paragraph()
    p.text = "944c25f386d8358b22f4f7734c67b12e55516bd2392ef54642803b93903a398"
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_BLUE
    p.space_before = Pt(4)

    p = tf.add_paragraph()
    p.text = "✔ Books balanced — debits = credits verified"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_before = Pt(16)

    # ==========================================
    # SLIDE 8: Measured Results (1,000 Incidents)
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s8)
    add_header(s8, "Measured Results: 1,000 Independent Incidents")

    # Confusion matrix metrics
    metrics = [
        ("119", "True Positives", ACCENT_GREEN),
        ("0", "False Positives", ACCENT_GREEN),
        ("49", "False Negatives", ACCENT_CORAL),
        ("832", "True Negatives", ACCENT_BLUE)
    ]
    for idx, (val, label, col) in enumerate(metrics):
        card = add_card(s8, Inches(0.8 + idx * 2.95), Inches(1.5), Inches(2.8), Inches(1.5))
        tb = s8.shapes.add_textbox(Inches(0.8 + idx * 2.95), Inches(1.6), Inches(2.8), Inches(1.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER

    add_card(s8, Inches(0.8), Inches(3.3), Inches(11.7), Inches(3.6), "HONEST ANALYSIS OF THE 49 FALSE NEGATIVES", ACCENT_YELLOW)
    tb = s8.shapes.add_textbox(Inches(1.1), Inches(4.0), Inches(11.1), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 100% PRECISION | 70.8% RECALL (N=1000, STABLE)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p = tf.add_paragraph()
    p.text = "• 43 of 49 Misses: Confidence stayed below the 70% activation floor on a genuine risk incident. The engine correctly held back rather than deploy capital on an uncertain signal."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "• 6 of 49 Misses: Occurred only after incident #965 of 1000, when cumulative deployment approached the 30% capital portfolio cap. The safety gate is an authentic constraint, not decorative."
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(10)

    # ==========================================
    # SLIDE 9: Real Trained Classifier on Harder Data
    # ==========================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s9)
    add_header(s9, "A Real Trained Classifier, Tested on Harder Data")

    add_card(s9, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tb = s9.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Trained on adversarial noisy dataset where bank failure is genuinely non-obvious:"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED

    p = tf.add_paragraph()
    p.text = "MODEL COMPARISON ON HELD-OUT TEST SET:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_before = Pt(14)

    p = tf.add_paragraph()
    p.text = "• Naive Rule:     27.9% Precision | 50.2% Recall | F1 = 0.359\n• XGBoost Model:  38.4% Precision | 54.5% Recall | F1 = 0.450"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = "Key Scientific Findings:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_YELLOW
    p.space_before = Pt(18)

    pts = [
        "Beats naive heuristics on both precision and recall simultaneously on held-out test data.",
        "Planted noise features correctly ranked dead last in feature importance (average rank 8.0 of 10).",
        "Threshold calibration: Calibrated optimal validation threshold at 0.225 (avoiding recall collapse at 0.70)."
    ]
    for pt in pts:
        p = tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 10: Found, Measured, and Fixed
    # ==========================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s10)
    add_header(s10, "Found, Measured, and Fixed: Engineering Iterations")

    add_card(s10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "CAPITAL ALLOCATION PROBLEM", ACCENT_CORAL)
    tb = s10.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The 30% Cap Ordering Flaw:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    c_pts = [
        "6 of 49 misses were caused by capital exhaustion near incident #965.",
        "Root cause: An ordering problem, not a confidence failure.",
        "FIX: Implemented Smallest-Transaction-First allocation + per-incident reserve cap.",
        "RESULT: 100% of capital-exhaustion misses rescued across 5 seeds with 0 new false positives."
    ]
    for pt in c_pts:
        p = tf.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(12)

    add_card(s10, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "AUTONOMOUS LEARNING LOOP", ACCENT_GREEN)
    tb = s10.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Simulating Real Drift & Retraining:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    l_pts = [
        "Simulated post-deployment feature drift in error streams.",
        "Triggered autonomous retrain on accumulated production telemetry.",
        "Verified on a third, completely untouched test batch.",
        "RESULT: F1 improved from 0.413 → 0.454 in 5 of 5 seed pairs tested."
    ]
    for pt in l_pts:
        p = tf.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(12)

    # ==========================================
    # SLIDE 11: Small, Honest, Reproducible
    # ==========================================
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s11)
    add_header(s11, "Financial Impact: Measured, Not Projected")

    fin_cards = [
        ("₹1.50 Cr", "PROTECTED VOLUME", ACCENT_BLUE),
        ("₹14,999.90", "FEE REVENUE @ 0.10%", ACCENT_GREEN),
        ("₹5,136.95", "NET FEE PROFIT", ACCENT_GREEN),
        ("₹29,999.80", "MERCHANT SAVINGS VS 0.30%", ACCENT_YELLOW)
    ]
    for idx, (val, lbl, col) in enumerate(fin_cards):
        add_card(s11, Inches(0.8 + idx * 2.95), Inches(1.5), Inches(2.8), Inches(1.8))
        tb = s11.shapes.add_textbox(Inches(0.8 + idx * 2.95), Inches(1.7), Inches(2.8), Inches(1.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER

    add_card(s11, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.2))
    tb = s11.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(11.0), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Literal stdout of `python scripts/batch_eval.py` — seed=42, n=1000"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p = tf.add_paragraph()
    p.text = "\"This isn't a new prediction model — Razorpay's is better than anything built in two weeks.\nIt's a product gap: turning a settlement product from reactive to predictive.\""
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(16)

    # ==========================================
    # SLIDE 12: Production-Ready & Deployed
    # ==========================================
    s12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s12)
    add_header(s12, "Production-Ready & Deployed (Full-Stack)")

    add_card(s12, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "LIVE SYSTEM INFRASTRUCTURE", ACCENT_BLUE)
    tb = s12.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    infra_pts = [
        "Interactive Console (Vercel): https://pisi-eosin.vercel.app",
        "Decision Engine (Render): https://settlement-guard-pisi-predictive-instant.onrender.com",
        "One-Click 'Simulate Downtime' Button: Triggers perception → inference → escalation → instant settlement.",
        "Webhook Receiver: 6 Razorpay event types with synchronous HMAC-SHA256 signature verification (<5ms latency).",
        "Uptime Monitoring: 24/7 keep-alive to eliminate cold starts."
    ]
    for pt in infra_pts:
        p = tf.paragraphs[0] if pt == infra_pts[0] else tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(12)

    add_card(s12, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "CAPITAL SAFEGUARDS & COMPLIANCE", ACCENT_GREEN)
    tb = s12.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    sec_pts = [
        "Tamper-Evident SHA-256 Audit Trail: Byte-for-byte verified with Web Crypto and Python hashlib.",
        "Strict Double-Entry Ledger: Balanced accounting with automatic T+2 replenishment reconciliation.",
        "3-Tier Escalation Matrix: HIGH (Auto-advance) / MEDIUM (Confirm) / LOW (Alert only).",
        "Hard Stopping Rules: 30% portfolio capital cap, 10 concurrent bridges, ₹50K per-tx limit."
    ]
    for pt in sec_pts:
        p = tf.paragraphs[0] if pt == sec_pts[0] else tf.add_paragraph()
        p.text = "✔ " + pt
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(12)

    # ==========================================
    # SLIDE 13: Beyond Original Scope
    # ==========================================
    s13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s13)
    add_header(s13, "What We've Built Beyond the Original Scope")

    items = [
        ("Real Razorpay API Integration", "Executes on-demand settlements with exact rupee-to-paise conversion"),
        ("Live Webhook Receiver", "Ingests real payment events, verifies HMAC signatures (<5ms latency)"),
        ("Trained XGBoost Classifier", "F1 = 0.6965 on core benchmark (and 0.450 on adversarial stress-test)"),
        ("Isolation Forest Anomaly Detection", "Catches novel out-of-distribution bank failures without labels"),
        ("Autonomous Learning Loop", "Auto-retrains on concept drift, validated across 5 seed pairs"),
        ("FastAPI Production Server", "Containerized and deployed on Render with Vercel reverse proxy")
    ]
    top_pos = 1.4
    for feat, imp in items:
        card = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1.2)

        tb = s13.shapes.add_textbox(Inches(1.1), Inches(top_pos + 0.1), Inches(3.8), Inches(0.65))
        p = tb.text_frame.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        tb2 = s13.shapes.add_textbox(Inches(5.0), Inches(top_pos + 0.1), Inches(7.2), Inches(0.65))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = imp
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        top_pos += 0.95

    # ==========================================
    # SLIDE 14: Closing
    # ==========================================
    s14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s14)

    tb = s14.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Smart Routing optimizes the path."
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "PISI optimizes the safety net."
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_before = Pt(16)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "PISI — Predictive Instant Settlement Intelligence"
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(40)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "One question: What if the merchant never had to ask?"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 15: Links & Resources
    # ==========================================
    s15 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s15)
    add_header(s15, "Links, Resources & Submission Details")

    add_card(s15, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tb = s15.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    links = [
        ("Interactive Live Dashboard:", "https://pisi-eosin.vercel.app", ACCENT_GREEN),
        ("Production REST API / Health:", "https://settlement-guard-pisi-predictive-instant.onrender.com/health", ACCENT_BLUE),
        ("GitHub Repository:", "https://github.com/asmitha2025/Settlement-Guard-PISI-Predictive-Instant-Settlement-Intelligence", ACCENT_BLUE),
        ("Target Track:", "Track 3 — AI Revenue Recovery (Razorpay AI Buildathon 2026)", TEXT_YELLOW),
        ("Key Innovation:", "Autonomous T+0 instant settlement pre-approval under 30% portfolio capital cap", TEXT_WHITE)
    ]

    for title, val, col in links:
        p = tf.paragraphs[0] if title == links[0][0] else tf.add_paragraph()
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_before = Pt(16)

        p2 = tf.add_paragraph()
        p2.text = "  " + val
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)

    # Save presentation
    output_path = "PISI_pitch_deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == '__main__':
    create_deck()
